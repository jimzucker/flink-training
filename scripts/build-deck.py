#!/usr/bin/env python3
"""Builds the final demo deck, in two versions, from one content model.

  docs/deck/final-demo.pptx           presenter's spine -- titles and
                                               evidence, nothing to read aloud
  docs/deck/final-demo-handout.pptx   stands alone -- prose on every
                                               slide, readable by someone who
                                               missed the demo

Both come from the CONTENT below so a number cannot drift between them. Every
figure here is quoted from a doc in this repo; the comment beside it says which.

  python3 -m venv .venv && .venv/bin/pip install python-pptx
  .venv/bin/python scripts/build-deck.py
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W, H = Inches(13.333), Inches(7.5)
DIAGRAM = "docs/deck/pipeline.png"

INK     = RGBColor(0x14, 0x22, 0x2E)   # near-black slate, all headings
BODY    = RGBColor(0x2E, 0x3D, 0x4A)
MUTED   = RGBColor(0x6B, 0x7A, 0x87)
ACCENT  = RGBColor(0x1B, 0x6B, 0x8A)   # deep teal -- the one accent
GOOD    = RGBColor(0x2C, 0x6E, 0x49)   # verified / passes
WARN    = RGBColor(0xA8, 0x4B, 0x2A)   # the ceiling, the caveat
RULE    = RGBColor(0xD8, 0xDF, 0xE4)
BAND    = RGBColor(0xF2, 0xF5, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

SANS, MONO = "Segoe UI", "Consolas"


def _tb(slide, x, y, w, h, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    return tf


def _run(p, text, size, color=BODY, bold=False, font=SANS, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.italic = Pt(size), bold, italic
    r.font.color.rgb, r.font.name = color, font
    return r


def _para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def heading(slide, title, kicker=None):
    """Kicker above, title below, hairline under. Every content slide."""
    y = Inches(0.52)
    if kicker:
        tf = _tb(slide, Inches(0.7), y, Inches(12), Inches(0.3))
        p = tf.paragraphs[0]
        r = _run(p, kicker.upper(), 11.5, ACCENT, bold=True)
        r.font._rPr.set("spc", "120")          # letter-spacing on the eyebrow
        y += Inches(0.34)
    tf = _tb(slide, Inches(0.7), y, Inches(12), Inches(0.8))
    _run(tf.paragraphs[0], title, 30, INK, bold=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7),
                                y + Inches(0.72), Inches(11.93), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    return y + Inches(1.05)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bullets(slide, top, items, size=17, gap=14):
    """items: (bold_lead, rest) tuples, or plain strings.

    The marker is a run in the paragraph rather than a drawn shape, so it stays
    on the line it belongs to no matter how the text wraps -- positioning a
    shape would mean predicting line heights, which drifts as soon as a bullet
    runs to two lines.
    """
    tf = _tb(slide, Inches(0.75), top, Inches(11.9), Inches(4.2))
    for i, item in enumerate(items):
        p = _para(tf, first=(i == 0))
        p.space_after = Pt(gap)
        _run(p, "\u25b8   ", size, ACCENT, bold=True)
        if isinstance(item, tuple):
            lead, rest = item
            _run(p, lead, size, INK, bold=True)
            _run(p, rest, size, BODY)
        else:
            _run(p, item, size, BODY)
    return tf


def table(slide, x, y, rows, widths, highlight=None, size=14, mono_from=1):
    """rows[0] is the header. highlight: set of row indexes to emphasise."""
    nr, nc = len(rows), len(rows[0])
    shape = slide.shapes.add_table(nr, nc, x, y, sum(widths), Inches(0.42) * nr)
    tbl = shape.table
    tbl.first_row = False
    for j, w in enumerate(widths):
        tbl.columns[j].width = w
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(0.44 if i == 0 else 0.40)
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = cell.margin_right = Inches(0.12)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = INK
            elif highlight and i in highlight:
                cell.fill.fore_color.rgb = BAND
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            emph = bool(highlight and i in highlight)
            _run(p, str(val), size,
                 WHITE if i == 0 else (INK if emph else BODY),
                 bold=(i == 0 or emph),
                 font=SANS if (i == 0 or j < mono_from) else MONO)
    return shape


def statline(slide, y, text, color=MUTED, size=13.5, italic=True):
    tf = _tb(slide, Inches(0.75), y, Inches(11.9), Inches(0.9))
    _run(tf.paragraphs[0], text, size, color, italic=italic)
    return tf


def prose(slide, y, text, size=15):
    tf = _tb(slide, Inches(0.75), y, Inches(11.9), Inches(1.6))
    _run(tf.paragraphs[0], text, size, BODY)
    return tf


# ------------------------------------------------------------------ slides --
# HANDOUT carries the prose; SPINE omits it. Numbers live here once.

def s_title(slide, handout):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(2.35))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK
    bar.line.fill.background(); bar.shadow.inherit = False

    tf = _tb(slide, Inches(0.9), Inches(0.72), Inches(11.5), Inches(1.4))
    p = tf.paragraphs[0]
    r = _run(p, "FINAL PROJECT", 12, RGBColor(0x7F, 0xB6, 0xCB), bold=True)
    r.font._rPr.set("spc", "160")
    p2 = tf.add_paragraph()
    _run(p2, "Positions and Market Value from Block Trades", 34, WHITE, bold=True)

    tf = _tb(slide, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.6))
    _run(tf.paragraphs[0],
         "A streaming pipeline that aggregates block trades two ways, joins prices, "
         "and publishes market value — with exactly-once delivery.", 16.5, BODY)

    rows = [["", ""],
            ["Runtime", "Apache Flink 1.20 · Java 17"],
            ["Transport", "Apache Kafka 3.9 (KRaft)"],
            ["Delivery", "exactly-once, verified by killing a task manager"],
            ["Observability", "Prometheus · Grafana"],
            ["Runs on", "one laptop, `docker compose up` — and on AWS unchanged"]]
    t = table(slide, Inches(0.9), Inches(3.65), rows[1:],
              [Inches(2.2), Inches(9.0)], size=13.5, mono_from=99)
    t.table.rows[0].height = Inches(0.42)

    tf = _tb(slide, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5))
    _run(tf.paragraphs[0], "The live demo is the deliverable. These slides are the handout and the backup.",
         13, MUTED, italic=True)
    notes(slide, "Open with the running stack already built but idle. The most persuasive "
                 "moment is graphs going from flat to flowing while people watch.")


def s_problem(slide, handout):
    top = heading(slide, "One order becomes four allocations", "The problem")
    bullets(slide, top, [
        ("A block trade is one order split across several accounts.  ",
         "A trade of 400 shares allocated to 4 accounts becomes 4 allocations of 100."),
        ("Publish positions two ways, in parallel.  ",
         "By symbol, and by account / sub-account / symbol."),
        ("Then join prices to those positions  ",
         "and publish market value the same two ways, once per key per window."),
    ])
    if handout:
        prose(slide, top + Inches(2.5),
              "That split is the whole reason the two aggregations differ. The account side sees four "
              "records for every one the symbol side sees, so sink 4 updates four times as often as "
              "sink 3 — while the quantities still reconcile, because the four allocations sum to the "
              "block. It is the first thing anyone notices on the dashboard and the first thing to "
              "have an answer ready for.")
    notes(slide, "Sink 4 updates 4x as often as sink 3. Quantities reconcile — the four "
                 "accounts sum to the block. The update counts differ by design.")


def s_design(slide, handout):
    """The project's own diagram, rendered by scripts/render-diagram.py."""
    top = heading(slide, "Six numbered elements, end to end", "The design")
    if not os.path.exists(DIAGRAM):
        raise SystemExit(f"{DIAGRAM} is missing -- run scripts/render-diagram.py first")
    with Image.open(DIAGRAM) as im:
        aspect = im.width / im.height
    w = Inches(12.2)
    h = Emu(int(w / aspect))
    slide.shapes.add_picture(DIAGRAM, Inches(0.57), top + Inches(0.28), width=w, height=h)
    statline(slide, top + Inches(0.28) + h + Inches(0.14),
             "Every edge carries its partition key and value. The diagram states the specified "
             "one-minute window; the demo overrides it to ten seconds, and the calculation is "
             "identical either way.")
    notes(slide, "Numbered left to right, in the order the demo talks through them. Point at 1 "
                 "and 2 as the inputs, 3-6 as the outputs to verify. If anyone reads '1 min' off "
                 "the diagram, that is the specification -- the demo runs a 10s window.")


def s_expected(slide, handout):
    top = heading(slide, "What the numbers should be, before we run anything",
                  "Expected output")

    # inputs and sinks sit side by side and end within half an inch of each
    # other; the demo's settings are a separate thing and get their own strip,
    # rather than being bolted onto the inputs table under a slash-compound
    # heading that has to cover both.
    table(slide, Inches(0.75), top,
          [["Input", "Rate"],
           ["Trades", "10 / sec"],
           ["Prices", "1 000 / sec"],
           ["Symbols", "4 unique"],
           ["Accounts", "4 unique"],
           ["Allocations per trade", "4 (one per account)"]],
          [Inches(3.1), Inches(2.3)], size=13.5)
    table(slide, Inches(6.35), top,
          [["#", "Sink", "Rate in the demo", "Keys"],
           ["3", "positions-by-symbol", "10 / sec", "4"],
           ["4", "positions-by-account", "40 / sec", "16"],
           ["5", "mv-by-symbol", "4 per 10 s", "4"],
           ["6", "mv-by-account", "16 per 10 s", "16"]],
          [Inches(0.45), Inches(3.05), Inches(1.9), Inches(0.9)], size=13.5)

    y = top + Inches(2.72)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.75), y, Inches(11.9), Inches(0.78))
    bar.fill.solid(); bar.fill.fore_color.rgb = BAND
    bar.line.color.rgb = RULE; bar.line.width = Pt(1)
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Inches(0.28)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, "Set for the demo:   ", 14, MUTED, bold=True)
    _run(p, "window ", 15, BODY)
    _run(p, "10 s", 15, INK, bold=True, font=MONO)
    _run(p, "  (specified: 1 min)     checkpoint interval ", 15, BODY)
    _run(p, "1 s", 15, INK, bold=True, font=MONO)
    _run(p, "     partitions ", 15, BODY)
    _run(p, "8", 15, INK, bold=True, font=MONO)

    statline(slide, y + Inches(1.02),
             "The window is 10 s so the market value sinks say something without a minute of "
             "waiting — the job defaults to the specified minute, and the calculation is identical "
             "either way. Committing to these numbers before the run is what makes the demo a "
             "verification rather than a tour; scripts/verify-topics.sh checks every one of them "
             "against the real topics, on every push.", GOOD, italic=False)
    notes(slide, "Say the numbers out loud BEFORE starting the generators. Then the dashboard "
                 "either matches or it does not. If asked why the window is 10s and not a minute: "
                 "only for the demo, so the sinks say something without a minute of waiting. The "
                 "job defaults to the specified minute and the verification runs against both.")


def s_live(slide, handout):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bar.fill.solid(); bar.fill.fore_color.rgb = INK
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = _tb(slide, Inches(1.1), Inches(2.6), Inches(11.2), Inches(2.0))
    p = tf.paragraphs[0]
    r = _run(p, "LIVE", 13, RGBColor(0x7F, 0xB6, 0xCB), bold=True)
    r.font._rPr.set("spc", "200")
    p2 = tf.add_paragraph()
    _run(p2, "Start the generators, then switch to Grafana", 34, WHITE, bold=True)
    p3 = tf.add_paragraph(); p3.space_before = Pt(16)
    _run(p3, "docker compose -f docker/compose.yml up -d --build     →     ./scripts/start-generators.sh",
         15, RGBColor(0xB9, 0xC6, 0xD0), font=MONO)
    notes(slide, "The generators come up idle on purpose, so the dashboard is visibly empty "
                 "until you start them. A stack that is already busy when it appears gives the "
                 "moment away.")


def s_numbers(slide, handout):
    top = heading(slide, "Every number on the screen has an answer", "Explain it")
    bullets(slide, top, [
        ("Why does sink 4 update four times as often as sink 3?  ",
         "Every block trade splits across four accounts. The quantities reconcile; the counts differ by design."),
        ("Why do the market value sinks look like steps?  ",
         "They emit once per key per window, not continuously."),
        ("Why is there one aborted checkpoint?  ",
         "At startup the first checkpoint triggers before every task is running. Harmless, and the panel says so."),
    ])
    if handout:
        prose(slide, top + Inches(2.6),
              "The requirement is to be able to explain any number on the dashboard, and if a "
              "question cannot be answered, to change the logging until it can. The aborted-checkpoint "
              "panel is an example of that going the other way first: it originally implied the "
              "guarantee was being retried, which would have raised a false alarm mid-demo.")
    notes(slide, "If you cannot answer something, say so and take it as an action item to change "
                 "the logging. That is the instruction in the requirements, and it beats an "
                 "explanation invented on the spot.")


def s_cases(slide, handout):
    top = heading(slide, "Both required cases pass", "Load")
    table(slide, Inches(0.75), top,
          [["", "orders/s asked", "prices/s", "orders through", "allocations", "order p50"],
           ["baseline", "10", "1 000", "8/s", "33/s", "519 ms"],
           ["case 1", "1 000", "1 000", "816/s", "3 269/s", "522 ms"],
           ["case 2", "10", "20 000", "8/s", "33/s", "513 ms"]],
          [Inches(1.5), Inches(2.2), Inches(1.5), Inches(2.0), Inches(1.7), Inches(1.5)],
          highlight={2, 3}, size=13.5)
    if handout:
        prose(slide, top + Inches(2.15),
              "Case 1 raised throughput a hundredfold and order latency did not move — the requirement "
              "allows latency to rise, and it did not need to. Case 2 raised the price rate twentyfold "
              "and latency was again unchanged, which settles the question left open when prices were "
              "made a broadcast: the concern was that every price would pass through the threads doing "
              "order work, and at twenty times the rate it does not.")
    else:
        statline(slide, top + Inches(2.15),
                 "A hundredfold on orders, twentyfold on prices — order latency unmoved in both.",
                 INK, 16, italic=False)
    notes(slide, "Case 2 is the one that settles the broadcast-prices question.")


def s_units(slide, handout):
    top = heading(slide, "Double the units, double the throughput", "Scaling")
    table(slide, Inches(0.75), top,
          [["units", "orders/sec", "vs previous", "Flink cores", "broker cores", "back-pressure"],
           ["2", "65 721", "—", "2.00 of 2", "0.25", "28.5%"],
           ["4", "129 056", "1.96×", "3.94 of 4", "0.39", "51.6%"]],
          [Inches(1.3), Inches(2.0), Inches(2.0), Inches(2.1), Inches(2.0), Inches(2.2)],
          highlight={2}, size=14)
    if handout:
        prose(slide, top + Inches(1.85),
              "A unit is one core and one degree of parallelism, bought together — which is exactly "
              "what a KPU is in Managed Service for Apache Flink. Two columns turn this from an "
              "assertion into a measurement. Flink used every core it was given, 2.00 of 2 and 3.94 "
              "of 4, so a unit bought is a unit worked. And the broker stayed under half a core, so "
              "Flink was unambiguously the constrained component — which is the precondition for "
              "showing that anything scales at all. Eight partitions were held constant across both "
              "cases, so nothing varied but the units.")
    else:
        statline(slide, top + Inches(1.85),
                 "A unit is one core and one degree of parallelism, bought together — a KPU.  "
                 "Flink used every core it was given; the broker stayed under half a core.",
                 INK, 16, italic=False)
    notes(slide, "You can only show that something scales when it is the thing that is constrained. "
                 "Step 10 got this wrong — it varied parallelism while Flink already had every core "
                 "it could use, and the answer came back flat.")


def s_units_chart(slide, handout):
    """The 2x result as a chart -- bars against a marked ideal, not a table.

    Two bars and a dashed line at exactly twice the first bar. The whole claim is
    whether the second bar reaches that line, so the line is drawn rather than
    described and the reader checks it in one glance.
    """
    top = heading(slide, "Two units, then four — against a perfect 2x", "Scaling")

    BASE = Inches(5.62)                     # baseline both bars stand on
    PLOT = Inches(3.42)                     # full height of the plot area
    FULL = 140000.0                         # value at the top of the plot
    def h(v):
        return Emu(int(PLOT * (v / FULL)))

    CASES = [(2, 65721, Inches(1.75)), (4, 129056, Inches(4.35))]
    BW = Inches(1.65)
    IDEAL = 65721 * 2                       # what a perfect 2x would reach

    # axis
    ax = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), BASE, Inches(5.9), Pt(1.25))
    ax.fill.solid(); ax.fill.fore_color.rgb = RULE
    ax.line.fill.background(); ax.shadow.inherit = False

    # the ideal-2x reference, drawn as a series of dashes so it reads as a target
    iy = BASE - h(IDEAL)
    x = Inches(1.30)
    while x < Inches(6.95):
        d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, iy, Inches(0.14), Pt(1.6))
        d.fill.solid(); d.fill.fore_color.rgb = WARN
        d.line.fill.background(); d.shadow.inherit = False
        x += Inches(0.26)
    tf = _tb(slide, Inches(5.55), iy - Inches(0.32), Inches(2.4), Inches(0.3))
    _run(tf.paragraphs[0], "perfect 2x  =  131,442", 11.5, WARN, bold=True)

    for units, val, bx in CASES:
        bh = h(val)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, BASE - bh, BW, bh)
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT if units == 4 else RGBColor(0x9F, 0xC2, 0xD1)
        bar.line.fill.background(); bar.shadow.inherit = False

        v = _tb(slide, bx - Inches(0.35), BASE - bh - Inches(0.46), BW + Inches(0.7),
                Inches(0.4), PP_ALIGN.CENTER)
        _run(v.paragraphs[0], f"{val:,}", 19, INK, bold=True)

        lab = _tb(slide, bx - Inches(0.35), BASE + Inches(0.10), BW + Inches(0.7),
                  Inches(0.66), PP_ALIGN.CENTER)
        p1 = lab.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        _run(p1, f"{units} units", 15, INK, bold=True)
        p2 = lab.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        _run(p2, f"{units:.2f} of {units} cores used" if units == 2 else "3.94 of 4 cores used",
             11.5, MUTED)

    # the multiplier, between the two bars
    m = _tb(slide, Inches(3.10), BASE - h(90000), Inches(1.5), Inches(0.8), PP_ALIGN.CENTER)
    p1 = m.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    _run(p1, "1.96x", 30, ACCENT, bold=True)
    p2 = m.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _run(p2, "orders / sec", 11, MUTED)

    # what the chart cannot show on its own
    rt = _tb(slide, Inches(7.75), top + Inches(0.30), Inches(4.9), Inches(4.0))
    for i, (lead, rest) in enumerate([
        ("A unit is one core and one degree of parallelism, bought together. ",
         "That is what a KPU is in Managed Service for Apache Flink."),
        ("Flink used every core it was given — ",
         "2.00 of 2, then 3.94 of 4. A unit bought is a unit worked."),
        ("The broker stayed under half a core. ",
         "Flink was the constrained component, and you can only show that something "
         "scales when it is the thing that is constrained."),
    ]):
        para = _para(rt, first=(i == 0)); para.space_after = Pt(15)
        _run(para, "\u25b8   ", 15, ACCENT, bold=True)
        _run(para, lead, 15, INK, bold=True)
        _run(para, rest, 15, BODY)

    if handout:
        statline(slide, Inches(6.42),
                 "Eight partitions were held constant across both cases, and the backlog was "
                 "drained with the producer stopped, so nothing varied but the units. The dashed "
                 "line is exactly twice the two-unit result; the four-unit bar reaches 98% of it.",
                 MUTED, 13)
    notes(slide, "The dashed line is exactly 2x the first bar. The whole claim is whether the "
                 "second bar reaches it -- 129,056 against 131,442, which is 98%. Step 10 got this "
                 "wrong by varying parallelism while Flink already had every core it could use, "
                 "and the answer came back flat.")


# Nine slides. The design comes before the problem -- the diagram makes the
# problem legible rather than the other way round. Scaling sits right after the
# live demo, while the pipeline is still on screen, instead of trailing the deck.
SLIDES = [s_title, s_design, s_problem, s_expected, s_live,
          s_units, s_units_chart, s_numbers, s_cases]


def build(path, handout):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in SLIDES:
        fn(blank(prs), handout)
    prs.save(path)
    return path


if __name__ == "__main__":
    out = "docs/deck"
    os.makedirs(out, exist_ok=True)
    import importlib.util                       # keeps the embedded PNG in step
    _spec = importlib.util.spec_from_file_location(
        "render_diagram", os.path.join(os.path.dirname(os.path.abspath(__file__)), "render-diagram.py"))
    _rd = importlib.util.module_from_spec(_spec); _spec.loader.exec_module(_rd)
    _rd.main()
    for name, handout in [("final-demo.pptx", False),
                          ("final-demo-handout.pptx", True)]:
        p = build(os.path.join(out, name), handout)
        print(f"  {p}  ({len(SLIDES)} slides, {'handout' if handout else 'spine'})")
