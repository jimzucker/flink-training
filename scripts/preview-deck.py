#!/usr/bin/env python3
"""Renders generated .pptx decks to an HTML page by reading their real geometry.

This is not a re-drawing of the content model -- it opens the .pptx files and
lays out whatever shapes are actually in them, at their actual coordinates, so
the page shows what the files contain rather than what they were meant to.

  .venv/bin/python scripts/preview-deck.py > /tmp/deck-preview.html
"""
import base64
import html
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

EMU_IN = 914400
SLIDE_W, SLIDE_H = 13.333 * EMU_IN, 7.5 * EMU_IN
PX_W = 1280.0                      # nominal stage width; everything scales off it

DECKS = [("spine", "Presenter's spine", "docs/deck/final-demo.pptx"),
         ("handout", "Standalone handout", "docs/deck/final-demo-handout.pptx")]


def pct(v, total):
    return f"{v / total * 100:.4f}%"


def fs(pt):
    """Point size -> container-relative units, so the stage scales cleanly."""
    return f"{float(pt) * (96.0 / 72.0) / PX_W * 100:.4f}cqw"


def rgb(color):
    try:
        if color and color.type is not None and color.rgb is not None:
            return "#%02X%02X%02X" % (color.rgb[0], color.rgb[1], color.rgb[2])
    except Exception:
        pass
    return None


ALIGN = {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right", PP_ALIGN.LEFT: "left"}


def runs_html(para, default_size=14.0):
    out = []
    for r in para.runs:
        st = []
        size = r.font.size.pt if r.font.size else default_size
        st.append(f"font-size:{fs(size)}")
        if r.font.bold:
            st.append("font-weight:700")
        if r.font.italic:
            st.append("font-style:italic")
        c = rgb(r.font.color)
        if c:
            st.append(f"color:{c}")
        name = r.font.name or ""
        if name.lower() in ("consolas", "courier new"):
            st.append("font-family:var(--mono)")
        txt = html.escape(r.text).replace(" ", "&nbsp;") if r.text.strip() == "" else html.escape(r.text)
        out.append(f'<span style="{";".join(st)}">{txt}</span>')
    return "".join(out) or "&nbsp;"


def textframe_html(tf):
    paras = []
    for p in tf.paragraphs:
        al = ALIGN.get(p.alignment, "left")
        before = p.space_before.pt if p.space_before else 0
        after = p.space_after.pt if p.space_after else 0
        style = (f"text-align:{al};margin-top:{fs(before)};margin-bottom:{fs(after)}")
        paras.append(f'<p style="{style}">{runs_html(p)}</p>')
    return "".join(paras)


def table_html(shape):
    tbl = shape.table
    widths = [c.width for c in tbl.columns]
    total = sum(widths) or 1
    rows = []
    for r_i, row in enumerate(tbl.rows):
        cells = []
        for c_i, cell in enumerate(row.cells):
            try:
                bg = rgb(cell.fill.fore_color)
            except (TypeError, AttributeError):
                bg = None
            st = [f"width:{widths[c_i] / total * 100:.3f}%"]
            if bg:
                st.append(f"background:{bg}")
            cells.append(f'<td style="{";".join(st)}">{textframe_html(cell.text_frame)}</td>')
        rows.append(f"<tr>{''.join(cells)}</tr>")
    return f'<table class="t">{"".join(rows)}</table>'


ROUND = {MSO_SHAPE.ROUNDED_RECTANGLE}
OVALS = {MSO_SHAPE.OVAL}
ARROWS = {MSO_SHAPE.RIGHT_ARROW}


def shape_html(sh):
    if None in (sh.left, sh.top, sh.width, sh.height):
        return ""
    pos = (f"left:{pct(sh.left, SLIDE_W)};top:{pct(sh.top, SLIDE_H)};"
           f"width:{pct(sh.width, SLIDE_W)};height:{pct(sh.height, SLIDE_H)}")

    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        img = sh.image
        b64 = base64.b64encode(img.blob).decode("ascii")
        return (f'<div class="sh" style="{pos}">'
                f'<img src="data:{img.content_type};base64,{b64}" alt=""></div>')

    if sh.has_table:
        return f'<div class="sh tblwrap" style="{pos}">{table_html(sh)}</div>'

    cls, extra = "sh", []
    auto = None
    if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            auto = sh.auto_shape_type
        except Exception:
            auto = None
        # a border-only shape has _NoFill, whose fore_color raises rather than
        # returning None -- asking for the type is not enough of a guard
        try:
            f = rgb(sh.fill.fore_color)
        except (TypeError, AttributeError):
            f = None
        if f:
            extra.append(f"background:{f}")
        try:
            ln = rgb(sh.line.color)
        except (TypeError, AttributeError):
            ln = None
        if ln and sh.line.width:
            extra.append(f"border:{max(1, round(sh.line.width.pt))}px solid {ln}")
        if auto in ROUND:
            extra.append("border-radius:0.9cqw")
        elif auto in OVALS:
            extra.append("border-radius:50%")
        elif auto in ARROWS:
            extra.append("clip-path:polygon(0 28%,72% 28%,72% 0,100% 50%,72% 100%,72% 72%,0 72%)")
    body = ""
    if sh.has_text_frame and sh.text_frame.text.strip():
        anchor = "center" if auto is not None else "flex-start"
        body = f'<div class="tf" style="justify-content:{anchor}">{textframe_html(sh.text_frame)}</div>'
    return f'<div class="{cls}" style="{pos};{";".join(extra)}">{body}</div>'


def deck_html(path):
    prs = Presentation(path)
    out = []
    for i, sl in enumerate(prs.slides, 1):
        shapes = "".join(shape_html(sh) for sh in sl.shapes)
        note = ""
        if sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip():
            note = (f'<div class="note"><span>Speaker note</span>'
                    f'{html.escape(sl.notes_slide.notes_text_frame.text.strip())}</div>')
        out.append(f'<figure class="slide"><div class="num">{i}</div>'
                   f'<div class="stage">{shapes}</div>{note}</figure>')
    return "\n".join(out)


TEMPLATE = """<title>Final Demo Deck</title>
<style>
  :root{
    --bg:#E9EDF0; --chrome:#FFFFFF; --ink:#16212B; --muted:#63727E;
    --line:#D3DBE1; --accent:#1B6B8A; --mono:"SF Mono",Menlo,Consolas,monospace;
    --sans:-apple-system,BlinkMacSystemFont,"Segoe UI",system-ui,sans-serif;
  }
  @media (prefers-color-scheme:dark){
    :root:not([data-theme="light"]){
      --bg:#11161A; --chrome:#1A2229; --ink:#E6ECF1; --muted:#95A5B2;
      --line:#2B363F; --accent:#63B4CE;
    }
  }
  :root[data-theme="dark"]{
    --bg:#11161A; --chrome:#1A2229; --ink:#E6ECF1; --muted:#95A5B2;
    --line:#2B363F; --accent:#63B4CE;
  }
  *{box-sizing:border-box}
  body{margin:0;background:var(--bg);color:var(--ink);font-family:var(--sans);
       -webkit-font-smoothing:antialiased}
  header{position:sticky;top:0;z-index:5;background:var(--chrome);
         border-bottom:1px solid var(--line);padding:14px 20px;
         display:flex;gap:18px;align-items:center;flex-wrap:wrap}
  h1{font-size:15px;margin:0;font-weight:650;letter-spacing:.01em}
  .sub{font-size:12.5px;color:var(--muted);margin:0}
  .tabs{margin-left:auto;display:flex;gap:6px}
  .tab{font:inherit;font-size:13px;padding:7px 13px;border-radius:7px;cursor:pointer;
       border:1px solid var(--line);background:transparent;color:var(--muted)}
  .tab.on{background:var(--accent);border-color:var(--accent);color:#fff;font-weight:600}
  main{max-width:1180px;margin:0 auto;padding:26px 20px 70px}
  .pane{display:none} .pane.on{display:block}
  .slide{margin:0 0 34px;position:relative}
  .num{font:600 11px/1 var(--mono);color:var(--muted);margin:0 0 7px 2px;
       letter-spacing:.08em}
  .stage{container-type:inline-size;position:relative;width:100%;aspect-ratio:1280/720;
         background:#fff;border:1px solid var(--line);border-radius:5px;overflow:hidden;
         box-shadow:0 1px 3px rgba(16,28,38,.10),0 8px 24px rgba(16,28,38,.06)}
  .sh{position:absolute;overflow:visible}
  .sh img{width:100%;height:100%;object-fit:contain;display:block}
  .tf{position:absolute;inset:0;display:flex;flex-direction:column;
      justify-content:flex-start;padding:.55cqw .7cqw}
  .tf p{margin:0;line-height:1.32;color:#2E3D4A}  /* slide-local: the stage is always white, like the real slide */
  .tblwrap{overflow:visible}
  .t{width:100%;height:100%;border-collapse:collapse;table-layout:fixed}
  .t td{vertical-align:middle;padding:.25cqw .55cqw;border:0}
  .t td p{line-height:1.25}
  .note{margin-top:9px;padding:10px 13px;background:var(--chrome);
        border:1px solid var(--line);border-left:3px solid var(--accent);
        border-radius:0 5px 5px 0;font-size:13px;line-height:1.5;color:var(--muted)}
  .note span{display:block;font-size:10.5px;font-weight:700;letter-spacing:.09em;
             text-transform:uppercase;color:var(--accent);margin-bottom:3px}
  @media (max-width:640px){ main{padding:16px 11px 50px} .tabs{margin-left:0;width:100%} }
</style>

<header>
  <div>
    <h1>Final demo deck — rendered from the .pptx files</h1>
    <p class="sub">Laid out from the real shape geometry in each file, not redrawn from the source.</p>
  </div>
  <div class="tabs">__TABS__</div>
</header>
<main>__PANES__</main>
<script>
  document.querySelectorAll('.tab').forEach(function(b){
    b.addEventListener('click', function(){
      document.querySelectorAll('.tab').forEach(function(x){x.classList.remove('on')});
      document.querySelectorAll('.pane').forEach(function(x){x.classList.remove('on')});
      b.classList.add('on');
      document.getElementById(b.dataset.k).classList.add('on');
      window.scrollTo({top:0});
    });
  });
</script>
"""


if __name__ == "__main__":
    panes = []
    tabs = []
    for idx, (key, label, path) in enumerate(DECKS):
        tabs.append(f'<button class="tab{" on" if idx == 0 else ""}" data-k="{key}">{label}</button>')
        panes.append(f'<section class="pane{" on" if idx == 0 else ""}" id="{key}">{deck_html(path)}</section>')
    sys.stdout.write(TEMPLATE.replace("__TABS__", "".join(tabs)).replace("__PANES__", "\n".join(panes)))
